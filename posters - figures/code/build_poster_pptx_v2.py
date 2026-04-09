#!/usr/bin/env python3
"""Build redesigned CeraWeek poster as editable PPTX.
36" x 24" (landscape), Stanford branding, cardinal red background.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Dimensions ──
W, H = Inches(36), Inches(24)
prs = Presentation()
prs.slide_width = W
prs.slide_height = H

slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

# ── Colors ──
CARDINAL   = RGBColor(0x8C, 0x15, 0x15)
DARK_CARD  = RGBColor(0x6B, 0x0F, 0x0F)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x33, 0x33, 0x33)
TEAL       = RGBColor(0x00, 0x6B, 0x5E)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY  = RGBColor(0x66, 0x66, 0x66)

# ── Background ──
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = CARDINAL

# ── Helper functions ──
def add_rect(left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color if fill_color else WHITE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.02
    return shape

def add_text_box(left, top, width, height, text, font_size=24, bold=False,
                 color=BLACK, alignment=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_text_box(left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_paragraph(tf, text, size=24, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                  space_after=Pt(4), font_name='Arial'):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    return p

def add_image(path, left, top, width=None, height=None):
    if width and height:
        pic = slide.shapes.add_picture(path, left, top, width, height)
    elif width:
        pic = slide.shapes.add_picture(path, left, top, width=width)
    else:
        pic = slide.shapes.add_picture(path, left, top)
    return pic

# ── Layout constants ──
MARGIN = Inches(0.3)
COL_GAP = Inches(0.25)
HEADER_H = Inches(2.2)
FOOTER_H = Inches(0.55)

content_top = HEADER_H + Inches(0.15)
content_bottom = H - FOOTER_H - Inches(0.1)
content_h = content_bottom - content_top

# 3 columns
col_w = (W - 2 * MARGIN - 2 * COL_GAP) / 3
col1_x = MARGIN
col2_x = MARGIN + col_w + COL_GAP
col3_x = MARGIN + 2 * (col_w + COL_GAP)

# ═══════════════════════════════════════════════════════
# HEADER
# ═══════════════════════════════════════════════════════
header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, HEADER_H)
header_bg.fill.solid()
header_bg.fill.fore_color.rgb = DARK_CARD
header_bg.line.fill.background()

# Gold accent line under header
gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, HEADER_H - Inches(0.06), W, Inches(0.06))
gold_line.fill.solid()
gold_line.fill.fore_color.rgb = RGBColor(0xD2, 0xC2, 0x95)
gold_line.line.fill.background()

add_text_box(Inches(1.5), Inches(0.15), Inches(33), Inches(0.9),
             'Direct Current (DC) Subtransmission Backbone for AI Factories',
             font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(Inches(1.5), Inches(0.95), Inches(33), Inches(0.5),
             'From Utility AC to 800 Vdc GPU Rails: Architecture, Efficiency, and Power Quality Analysis',
             font_size=26, bold=False, color=RGBColor(0xDD, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)

add_text_box(Inches(1.5), Inches(1.4), Inches(33), Inches(0.4),
             'Jane Yang (yjane@stanford.edu)    |    Liang Min (liangmin@stanford.edu)',
             font_size=26, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(Inches(1.5), Inches(1.75), Inches(33), Inches(0.3),
             'Stanford University    |    CeraWeek 2026',
             font_size=22, bold=False, color=RGBColor(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# COLUMN 1: ABSTRACT & MOTIVATION + PROPOSED ARCHITECTURE
# ═══════════════════════════════════════════════════════

# --- Section 1: Abstract & Motivation ---
sec1_h = Inches(10.2)
add_rect(col1_x, content_top, col_w, sec1_h, WHITE)

tb1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col1_x, content_top, col_w, Inches(0.5))
tb1.fill.solid(); tb1.fill.fore_color.rgb = CARDINAL; tb1.line.fill.background()
add_text_box(col1_x + Inches(0.1), content_top + Inches(0.03), col_w - Inches(0.2), Inches(0.44),
             'ABSTRACT & MOTIVATION', font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

tf1 = add_rich_text_box(col1_x + Inches(0.2), content_top + Inches(0.6),
                         col_w - Inches(0.4), Inches(4.2))

add_paragraph(tf1, 'Why Revisit the Power Path Now?',
              size=26, bold=True, color=CARDINAL, space_after=Pt(6))

add_paragraph(tf1, 'AI-training loads are changing the problem from a purely steady-state '
              'efficiency problem into an architecture + power-dynamics problem. '
              'GPU is the dominant power consumer (~65%), so synchronized GPU activity '
              'directly shows up at the server/rack/datacenter level.',
              size=19, bold=False, color=BLACK, space_after=Pt(8))

add_paragraph(tf1, 'Three Scenarios in the Industry:',
              size=22, bold=True, color=TEAL, space_after=Pt(4))

add_paragraph(tf1, '1. Traditional AC-centric: Utility AC through 5-6 conversion stages '
              '(transformer \u2192 UPS \u2192 PDU \u2192 PSU \u2192 VRM) to GPU. Mature but '
              'inefficient with repeated AC-DC-AC conversions.',
              size=18, bold=False, color=BLACK, space_after=Pt(4))

add_paragraph(tf1, '2. AC-fed SST / 800 Vdc pod (NVIDIA et al.): Grid-side AC/DC feeds '
              'HF isolated pods producing 800 Vdc. Reduces stages but AC still enters '
              'upstream \u2014 harmonics and PF issues persist at each pod.',
              size=18, bold=False, color=BLACK, space_after=Pt(4))

add_paragraph(tf1, '3. Proposed MVDC hub + isolated DC pod: Single MV AC/DC front-end '
              'feeds an MVDC backbone. Harmonics handled once. DC-native path for '
              'storage, solar, and rack buffering. Fewest conversion stages.',
              size=18, bold=False, color=BLACK, space_after=Pt(6))

add_paragraph(tf1, 'Key Challenges with AC Distribution:',
              size=22, bold=True, color=RGBColor(0xC6, 0x28, 0x28), space_after=Pt(4))

add_paragraph(tf1, '\u2022 Harmonic distortion from thousands of nonlinear server PSUs\n'
              '\u2022 Voltage regulation across distributed AC buses\n'
              '\u2022 Reactive power compensation at every conversion stage\n'
              '\u2022 Synchronous GPU load swings (0.1\u201320 Hz) stress AC infrastructure',
              size=18, bold=False, color=BLACK, space_after=Pt(4))

# GPU dynamics figure
gpu_fig_top = content_top + Inches(6.6)
add_image('/home/ubuntu/poster/v2_fig_gpu_dynamics.png',
          col1_x + Inches(0.15), gpu_fig_top,
          width=col_w - Inches(0.3))

# Scale / Dynamics / Spectrum
tf_motiv = add_rich_text_box(col1_x + Inches(0.2), gpu_fig_top + Inches(2.2),
                              col_w - Inches(0.4), Inches(1.2))
add_paragraph(tf_motiv, '\u2022 Scale: Deployments growing from <10 MW to >100 MW; single training '
              'jobs span >100,000 GPUs.',
              size=16, bold=False, color=BLACK, space_after=Pt(2))
add_paragraph(tf_motiv, '\u2022 Dynamics: Synchronous compute/communication phases create periodic '
              'power swings visible beyond node and rack.',
              size=16, bold=False, color=BLACK, space_after=Pt(2))
add_paragraph(tf_motiv, '\u2022 Spectrum: FFT energy concentrates around 0.2\u20133 Hz inside a broader '
              '0.1\u201320 Hz concern band.',
              size=16, bold=False, color=BLACK, space_after=Pt(2))

# --- Section 2: Proposed Architecture ---
sec2_top = content_top + sec1_h + Inches(0.2)
sec2_h = content_h - sec1_h - Inches(0.2)
add_rect(col1_x, sec2_top, col_w, sec2_h, WHITE)

tb2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col1_x, sec2_top, col_w, Inches(0.5))
tb2.fill.solid(); tb2.fill.fore_color.rgb = CARDINAL; tb2.line.fill.background()
add_text_box(col1_x + Inches(0.1), sec2_top + Inches(0.03), col_w - Inches(0.2), Inches(0.44),
             'PROPOSED ARCHITECTURE: THREE SCENARIOS', font_size=26, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Architecture diagram
arch_img_top = sec2_top + Inches(0.6)
arch_img_h = sec2_h - Inches(3.3)
add_image('/home/ubuntu/poster/v2_fig_architecture.png',
          col1_x + Inches(0.1), arch_img_top,
          width=col_w - Inches(0.2))

# Conceptual target stack text
tf_arch = add_rich_text_box(col1_x + Inches(0.2), sec2_top + sec2_h - Inches(2.8),
                             col_w - Inches(0.4), Inches(2.6))
add_paragraph(tf_arch, 'Conceptual Target Stack:',
              size=22, bold=True, color=TEAL, space_after=Pt(4))
add_paragraph(tf_arch, 'Utility AC \u2192 MV AC/DC \u2192 MVDC Hub \u2192 Isolated DC Pod \u2192 GPU',
              size=20, bold=True, color=BLACK, space_after=Pt(6))
add_paragraph(tf_arch, 'The claim: the pod gets cleaner while the facility-level interface '
              'becomes more centralized and more strategic. This is the architecture to '
              'test, not to assume.',
              size=18, bold=False, color=DARK_GRAY, space_after=Pt(6))
add_paragraph(tf_arch, 'Efficiency formula:  \u03b7_total = \u220f(\u03b7_i) for each stage i\n'
              'AC: N=5\u20136, \u03b7_i \u2248 0.92\u20130.98 \u21d2 \u03b7_total \u2248 81%\n'
              'DC: N=3\u20134, \u03b7_i \u2248 0.97\u20130.995 \u21d2 \u03b7_total \u2248 94.6%',
              size=17, bold=False, color=BLACK, space_after=Pt(4))
add_paragraph(tf_arch, 'Source: MVDC Position Paper, Stanford Energy, 2025',
              size=14, bold=False, color=DARK_GRAY, space_after=Pt(2))


# ═══════════════════════════════════════════════════════
# COLUMN 2: THREE BENEFITS
# ═══════════════════════════════════════════════════════

# --- Benefit 1: Cumulative Efficiency ---
ben1_h = Inches(7.0)
add_rect(col2_x, content_top, col_w, ben1_h, WHITE)

tb3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col2_x, content_top, col_w, Inches(0.5))
tb3.fill.solid(); tb3.fill.fore_color.rgb = CARDINAL; tb3.line.fill.background()
add_text_box(col2_x + Inches(0.1), content_top + Inches(0.03), col_w - Inches(0.2), Inches(0.44),
             'BENEFIT 1: CUMULATIVE EFFICIENCY', font_size=26, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

tf_eff = add_rich_text_box(col2_x + Inches(0.2), content_top + Inches(0.6),
                            col_w - Inches(0.4), Inches(2.2))
add_paragraph(tf_eff, 'Each AC-DC or DC-DC conversion stage incurs 1\u20135% loss. The MVDC '
              'architecture eliminates redundant conversions, achieving 94.6% end-to-end '
              'efficiency vs. 81.1% for traditional AC \u2014 a 13.5 percentage-point gain.',
              size=19, bold=False, color=BLACK, space_after=Pt(6))

add_paragraph(tf_eff, 'For a 100 MW AI factory, this translates to ~13.5 MW saved, or '
              'approximately $11.8M/year in electricity costs at $0.10/kWh.',
              size=19, bold=True, color=TEAL, space_after=Pt(6))

add_paragraph(tf_eff, '\u0394\u03b7 = \u03b7_DC \u2212 \u03b7_AC = 94.6% \u2212 81.1% = 13.5%',
              size=18, bold=True, color=RGBColor(0x00, 0x6B, 0x5E), space_after=Pt(2))

# Efficiency chart
eff_img_top = content_top + Inches(3.3)
add_image('/home/ubuntu/poster/v2_fig_efficiency.png',
          col2_x + Inches(0.1), eff_img_top,
          width=col_w - Inches(0.2))

# Source citation
add_text_box(col2_x + Inches(0.2), content_top + Inches(6.3), col_w - Inches(0.4), Inches(0.5),
             'Sources: IEEE PELS 2024; EPRI DC Data Center Report; Rasmussen WP#63',
             font_size=14, bold=False, color=DARK_GRAY)

# --- Benefit 2: Harmonics & Power Quality ---
ben2_top = content_top + ben1_h + Inches(0.2)
ben2_h = Inches(7.0)
add_rect(col2_x, ben2_top, col_w, ben2_h, WHITE)

tb4 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col2_x, ben2_top, col_w, Inches(0.5))
tb4.fill.solid(); tb4.fill.fore_color.rgb = CARDINAL; tb4.line.fill.background()
add_text_box(col2_x + Inches(0.1), ben2_top + Inches(0.03), col_w - Inches(0.2), Inches(0.44),
             'BENEFIT 2: HARMONICS & POWER QUALITY', font_size=26, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

tf_harm = add_rich_text_box(col2_x + Inches(0.2), ben2_top + Inches(0.6),
                             col_w - Inches(0.4), Inches(2.4))
add_paragraph(tf_harm, 'In AC-centric architectures, every server PSU injects harmonics '
              'and requires power factor correction (PFC). With thousands of servers, '
              'harmonic management is distributed across 4\u20135 points.',
              size=19, bold=False, color=BLACK, space_after=Pt(6))

add_paragraph(tf_harm, 'With the MVDC hub, harmonic and PF handling become centralized '
              'once at the MV AC/DC front-end, instead of being repeatedly embedded '
              'in every downstream pod.',
              size=19, bold=True, color=TEAL, space_after=Pt(6))

add_paragraph(tf_harm, 'For utilities: AC-side power quality needs to be managed only once '
              'at the grid interface \u2014 dramatically simplifying compliance with IEEE 519 '
              'harmonic limits and reducing equipment cost.',
              size=19, bold=False, color=BLACK, space_after=Pt(2))

# Harmonics figure
harm_img_top = ben2_top + Inches(3.5)
add_image('/home/ubuntu/poster/v2_fig_harmonics.png',
          col2_x + Inches(0.1), harm_img_top,
          width=col_w - Inches(0.2))

# --- Benefit 3: Voltage Management & DC-Native Integration ---
ben3_top = ben2_top + ben2_h + Inches(0.2)
ben3_h = content_h - ben1_h - ben2_h - Inches(0.4)
add_rect(col2_x, ben3_top, col_w, ben3_h, WHITE)

tb5 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col2_x, ben3_top, col_w, Inches(0.5))
tb5.fill.solid(); tb5.fill.fore_color.rgb = CARDINAL; tb5.line.fill.background()
add_text_box(col2_x + Inches(0.1), ben3_top + Inches(0.03), col_w - Inches(0.2), Inches(0.44),
             'BENEFIT 3: VOLTAGE MGMT & DC-NATIVE PATH', font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

tf_volt = add_rich_text_box(col2_x + Inches(0.2), ben3_top + Inches(0.6),
                             col_w - Inches(0.4), Inches(2.2))
add_paragraph(tf_volt, 'The MVDC backbone provides a cleaner DC-native path for rack '
              'buffering, battery storage (BESS), solar PV, and future modular AI '
              'power blocks \u2014 all connecting directly to the DC bus without AC-DC-AC '
              're-conversion.',
              size=19, bold=False, color=BLACK, space_after=Pt(6))

add_paragraph(tf_volt, 'Voltage regulation on a DC bus is inherently simpler: no reactive '
              'power, no frequency control, no phase balancing. A single DC voltage '
              'setpoint replaces complex multi-variable AC control.',
              size=19, bold=False, color=BLACK, space_after=Pt(6))

add_paragraph(tf_volt, 'P_DC = 2\u00b7V\u00b7I  vs.  P_AC = \u221a3\u00b7V\u00b7I\u00b7cos\u03c6\n'
              'DC delivers ~15% more power per conductor pair.',
              size=18, bold=True, color=TEAL, space_after=Pt(2))

# Voltage management figure
volt_img_top = ben3_top + Inches(3.2)
add_image('/home/ubuntu/poster/v2_fig_voltage_mgmt.png',
          col2_x + Inches(0.1), volt_img_top,
          width=col_w - Inches(0.2))


# ═══════════════════════════════════════════════════════
# COLUMN 3: COMPARISON TABLE + CONCLUSION
# ═══════════════════════════════════════════════════════

# --- Comparison Table ---
tbl_sec_h = Inches(12.5)
add_rect(col3_x, content_top, col_w, tbl_sec_h, WHITE)

tb6 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col3_x, content_top, col_w, Inches(0.5))
tb6.fill.solid(); tb6.fill.fore_color.rgb = CARDINAL; tb6.line.fill.background()
add_text_box(col3_x + Inches(0.1), content_top + Inches(0.03), col_w - Inches(0.2), Inches(0.44),
             'BENEFITS & DRAWBACKS BY ARCHITECTURE', font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(col3_x + Inches(0.2), content_top + Inches(0.55), col_w - Inches(0.4), Inches(0.5),
             'The proposed path is not "free efficiency"; it exchanges downstream simplicity '
             'for a tougher MVDC infrastructure problem.',
             font_size=16, bold=False, color=DARK_GRAY)

# Build comparison table
tbl_top = content_top + Inches(1.15)
tbl_left = col3_x + Inches(0.15)
tbl_width = col_w - Inches(0.3)

rows = 8
cols = 4
table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_width, Inches(9.2))
table = table_shape.table

col_widths = [Inches(1.8), Inches(2.8), Inches(2.8), Inches(2.8)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

headers = ['Comparison\nAxis', 'Traditional\nAC-centric', 'AC-fed SST /\n800 Vdc', 'MVDC hub +\nisolated pod']
header_colors = [DARK_GRAY, RGBColor(0x4A, 0x4A, 0x4A), CARDINAL, TEAL]

for j, (hdr, hcolor) in enumerate(zip(headers, header_colors)):
    cell = table.cell(0, j)
    cell.text = hdr
    cell.fill.solid()
    cell.fill.fore_color.rgb = hcolor
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

data = [
    ['Facility-side\ninterface',
     'AC-native utility +\nfacility protection',
     'AC input still present\nat pod / cluster front end',
     'Centralized MV AC/DC\nfront end feeding\nan MVDC bus'],
    ['Local pod\ncomplexity',
     'No SST pod, but repeated\ndownstream AC/DC + DC/DC',
     'High local sophistication:\nPFC + HF isolated link\n+ 800 Vdc output',
     'Simpler pod: HF isolated\nlink remains, grid-facing\nAC/DC moves upstream'],
    ['DC-native\nintegration',
     'Extra interfaces for\nstorage / solar / other\nDC systems',
     'Better local DC handoff,\nbut AC still enters\nupstream',
     'Best native fit for DC\nbuffering, storage, solar,\nand rack-level DC'],
    ['Harmonic / PF\nownership',
     'Distributed across AC\nfacility equipment and\nserver-side interfaces',
     'Handled locally at each\npod or cluster interface',
     'More centralized at the\ncommon MV front end'],
    ['Maturity /\nstandards',
     'Highest',
     'Prototype to emerging\nproduct concepts',
     'Lowest today; protection\nand code pathway remain\nmajor hurdles'],
    ['Main\nbottleneck',
     'Repeated conversion\nlayers, copper, and\nfacility sprawl',
     'SiC cost, HFT insulation,\ncapacitor stress, control',
     'MVDC protection, central\navailability, and safe\nfault handling'],
]

for i, row_data in enumerate(data):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i + 1, j)
        cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = (j == 0)
            p.font.color.rgb = BLACK
            p.font.name = 'Arial'
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if j == 3:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
        elif j == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

# Bottom line row
bottom_cell = table.cell(7, 0)
bottom_cell_end = table.cell(7, 3)
bottom_cell.merge(bottom_cell_end)
bottom_cell.text = ('Bottom line: the MVDC-hub proposal is strongest as a system-architecture story, '
                    'not as a claim that one isolated converter magically solves the whole facility by itself.')
for p in bottom_cell.text_frame.paragraphs:
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = CARDINAL
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
bottom_cell.fill.solid()
bottom_cell.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)

# Key metrics summary box
metrics_top = content_top + Inches(10.6)
metrics_bg = add_rect(col3_x + Inches(0.15), metrics_top,
                       col_w - Inches(0.3), Inches(1.7),
                       RGBColor(0xE8, 0xF5, 0xE9), TEAL, Pt(2))

tf_metrics = add_rich_text_box(col3_x + Inches(0.3), metrics_top + Inches(0.1),
                                col_w - Inches(0.6), Inches(1.5))
add_paragraph(tf_metrics, 'Key Metrics Summary (100 MW AI Factory):',
              size=20, bold=True, color=TEAL, space_after=Pt(4))
add_paragraph(tf_metrics, '\u2022 End-to-end efficiency: 94.6% (MVDC) vs. 81.1% (AC) \u2014 \u0394\u03b7 = 13.5%\n'
              '\u2022 Power saved: ~13.5 MW \u2192 ~$11.8M/year at $0.10/kWh\n'
              '\u2022 Conversion stages: 3\u20134 (MVDC) vs. 5\u20136 (AC)\n'
              '\u2022 PQ management points: 1 (MVDC) vs. 4\u20135 (AC)\n'
              '\u2022 DC-native integration: BESS, solar, rack buffering \u2014 no re-conversion',
              size=16, bold=False, color=BLACK, space_after=Pt(2))

# --- Conclusion & Future Work ---
conc_top = content_top + tbl_sec_h + Inches(0.2)
conc_h = content_h - tbl_sec_h - Inches(0.2)
add_rect(col3_x, conc_top, col_w, conc_h, WHITE)

tb7 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col3_x, conc_top, col_w, Inches(0.5))
tb7.fill.solid(); tb7.fill.fore_color.rgb = CARDINAL; tb7.line.fill.background()
add_text_box(col3_x + Inches(0.1), conc_top + Inches(0.03), col_w - Inches(0.2), Inches(0.44),
             'CONCLUSION & FUTURE WORK', font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

tf_conc = add_rich_text_box(col3_x + Inches(0.2), conc_top + Inches(0.6),
                             col_w - Inches(0.4), conc_h - Inches(0.7))

add_paragraph(tf_conc, 'Conclusion',
              size=24, bold=True, color=TEAL, space_after=Pt(4))
add_paragraph(tf_conc, 'The MVDC subtransmission backbone offers a compelling system-architecture '
              'path for next-generation AI factories: fewer conversion stages, centralized '
              'power quality management, simpler voltage regulation, and native DC integration '
              'for storage and renewables.',
              size=18, bold=False, color=BLACK, space_after=Pt(8))

add_paragraph(tf_conc, 'Suggested Evaluation Plan:',
              size=22, bold=True, color=CARDINAL, space_after=Pt(4))
add_paragraph(tf_conc, '\u2022 Compare conversion count, cable section / ohmic loss, and equipment stack\n'
              '\u2022 Map fault-clearing location and protection responsibility\n'
              '\u2022 Quantify SiC, transformer, and capacitor stress / BOM concentration\n'
              '\u2022 Decide where buffering belongs: GPU, rack, pod, or MVDC hub',
              size=17, bold=False, color=BLACK, space_after=Pt(8))

add_paragraph(tf_conc, 'Future Work:',
              size=22, bold=True, color=CARDINAL, space_after=Pt(4))
add_paragraph(tf_conc, '\u2022 MVDC protection schemes and fault isolation strategies\n'
              '\u2022 Techno-economic analysis at 100\u2013500 MW scale\n'
              '\u2022 Hardware-in-the-loop validation of MVDC hub + isolated pod\n'
              '\u2022 Standards development for DC data center distribution',
              size=17, bold=False, color=BLACK, space_after=Pt(6))

add_paragraph(tf_conc, 'References',
              size=20, bold=True, color=CARDINAL, space_after=Pt(3))
add_paragraph(tf_conc, '[1] IEA, "Data Centres and AI," 2024\n'
              '[2] IEEE PELS, DC Distribution for Data Centers, 2024\n'
              '[3] EPRI, DC Power for Data Centers Report, 2023\n'
              '[4] NVIDIA, GB200 Power Architecture Whitepaper, 2024\n'
              '[5] Rasmussen, "AC vs DC Power Distribution," APC WP#63\n'
              '[6] Pratt et al., "Evaluation of 380V DC," INTELEC 2007\n'
              '[7] AlLee & Tschudi, "Edison Redux," IEEE Power & Energy, 2012',
              size=14, bold=False, color=DARK_GRAY, space_after=Pt(2))


# ═══════════════════════════════════════════════════════
# FOOTER
# ═══════════════════════════════════════════════════════
footer_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - FOOTER_H, W, FOOTER_H)
footer_bg.fill.solid()
footer_bg.fill.fore_color.rgb = DARK_CARD
footer_bg.line.fill.background()

add_text_box(Inches(0.5), H - FOOTER_H + Inches(0.1), Inches(14), Inches(0.35),
             'Jane Yang (yjane@stanford.edu)  |  Liang Min (liangmin@stanford.edu)  |  Stanford University',
             font_size=20, bold=True, color=WHITE)

add_text_box(Inches(22), H - FOOTER_H + Inches(0.1), Inches(13.5), Inches(0.35),
             'CeraWeek 2026  |  DC Subtransmission Backbone for AI Factories',
             font_size=20, bold=False, color=RGBColor(0xCC, 0xCC, 0xCC),
             alignment=PP_ALIGN.RIGHT)

# ── Save ──
out_path = '/home/ubuntu/poster/poster_v2.pptx'
prs.save(out_path)
print(f"Poster saved to {out_path}")
print(f"File size: {os.path.getsize(out_path) / 1024:.0f} KB")
