#!/usr/bin/env python3
"""Build CeraWeek poster v3: reduced text, bigger fonts, balanced, academic Stanford theme.
36" x 24" landscape, editable PPTX.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import json
import os

BASE_DIR = Path(__file__).resolve().parents[2]
FIG_DIR = Path(__file__).resolve().parents[1] / 'figures'
REPORT_PATH = BASE_DIR / 'source_backed_model_report.json'
OUTPUT_PATH = Path(__file__).resolve().parents[1] / 'poster_v3.pptx'

report = json.loads(REPORT_PATH.read_text(encoding='utf-8'))
results = {item['name']: item for item in report['results']}
scenario_1 = results['traditional_ac']
scenario_2 = results['ac_fed_sst_800vdc']
scenario_3 = results['proposed_mvdc_backbone']

eff_s1 = scenario_1['full_load_total_efficiency'] * 100.0
eff_s2 = scenario_2['full_load_total_efficiency'] * 100.0
eff_s3 = scenario_3['full_load_total_efficiency'] * 100.0
loss_gwh_s1 = scenario_1['annual_loss_mwh'] / 1000.0
loss_gwh_s2 = scenario_2['annual_loss_mwh'] / 1000.0
loss_gwh_s3 = scenario_3['annual_loss_mwh'] / 1000.0
cost_musd_s1 = scenario_1['annual_loss_cost_usd'] / 1e6
cost_musd_s2 = scenario_2['annual_loss_cost_usd'] / 1e6
cost_musd_s3 = scenario_3['annual_loss_cost_usd'] / 1e6

delta_eff_s3_vs_s1 = eff_s3 - eff_s1
delta_loss_gwh_s3_vs_s1 = loss_gwh_s1 - loss_gwh_s3
delta_cost_musd_s3_vs_s1 = cost_musd_s1 - cost_musd_s3
delta_input_mw_s3_vs_s1 = scenario_1['full_load_input_mw'] - scenario_3['full_load_input_mw']
pq_points_s1 = scenario_1["innovation_metrics"]["ac_harmonic_injection_points"]
pq_points_s2 = scenario_2["innovation_metrics"]["ac_harmonic_injection_points"]
pq_points_s3 = scenario_3["innovation_metrics"]["ac_harmonic_injection_points"]
major_conv_s1 = scenario_1["innovation_metrics"]["major_conversion_stages"]
major_conv_s2 = scenario_2["innovation_metrics"]["major_conversion_stages"]
major_conv_s3 = scenario_3["innovation_metrics"]["major_conversion_stages"]

opendss_s1 = scenario_1.get("opendss_validation", {}).get("base_snapshot", {})
opendss_s2 = scenario_2.get("opendss_validation", {}).get("base_snapshot", {})
opendss_s3 = scenario_3.get("opendss_validation", {}).get("base_snapshot", {})

W, H = Inches(36), Inches(24)
prs = Presentation()
prs.slide_width = W
prs.slide_height = H
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── Colors ──
CARDINAL   = RGBColor(0x8C, 0x15, 0x15)
DARK_CARD  = RGBColor(0x6B, 0x0F, 0x0F)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x33, 0x33, 0x33)
TEAL       = RGBColor(0x00, 0x6B, 0x5E)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY  = RGBColor(0x66, 0x66, 0x66)
GOLD       = RGBColor(0xD2, 0xC2, 0x95)

bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = CARDINAL

# ── Helpers ──
def add_rect(l, t, w, h, fc=WHITE, bc=None, bw=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if bc: s.line.color.rgb = bc; s.line.width = bw
    else: s.line.fill.background()
    s.adjustments[0] = 0.015
    return s

def add_box(l, t, w, h, fc=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc; s.line.fill.background()
    return s

def tb(l, t, w, h, text, sz=24, bold=False, color=BLACK, align=PP_ALIGN.LEFT, fn='Arial'):
    txB = slide.shapes.add_textbox(l, t, w, h)
    tf = txB.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.bold = bold; p.font.color.rgb = color; p.font.name = fn; p.alignment = align
    return txB

def rtb(l, t, w, h):
    txB = slide.shapes.add_textbox(l, t, w, h)
    tf = txB.text_frame; tf.word_wrap = True; return tf

def ap(tf, text, sz=24, bold=False, color=BLACK, align=PP_ALIGN.LEFT, sa=Pt(4), fn='Arial'):
    if len(tf.paragraphs)==1 and tf.paragraphs[0].text=='': p=tf.paragraphs[0]
    else: p=tf.add_paragraph()
    p.text=text; p.font.size=Pt(sz); p.font.bold=bold; p.font.color.rgb=color
    p.font.name=fn; p.alignment=align; p.space_after=sa; return p

def img(path, l, t, w=None, h=None, width=None, height=None):
    w = w or width; h = h or height
    if w and h: return slide.shapes.add_picture(path, l, t, w, h)
    elif w: return slide.shapes.add_picture(path, l, t, width=w)
    else: return slide.shapes.add_picture(path, l, t)

def section_title(l, t, w, text, sz=28):
    add_box(l, t, w, Inches(0.55), CARDINAL)
    tb(l+Inches(0.1), t+Inches(0.05), w-Inches(0.2), Inches(0.45),
       text, sz=sz, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Layout ──
M = Inches(0.3)
GAP = Inches(0.25)
HDR_H = Inches(2.1)
FTR_H = Inches(0.5)
CT = HDR_H + Inches(0.15)
CB = H - FTR_H - Inches(0.1)
CH = CB - CT
CW = (W - 2*M - 2*GAP) / 3
C1 = M; C2 = M + CW + GAP; C3 = M + 2*(CW+GAP)

# ═══════════════════════════════════════
# HEADER
# ═══════════════════════════════════════
add_box(0, 0, W, HDR_H, DARK_CARD)
add_box(0, HDR_H - Inches(0.05), W, Inches(0.05), GOLD)

tb(Inches(1.5), Inches(0.12), Inches(33), Inches(0.85),
   'Direct Current (DC) Subtransmission Backbone for AI Factories',
   sz=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tb(Inches(1.5), Inches(0.88), Inches(33), Inches(0.45),
   'From Utility AC to 800 Vdc GPU Rails: Architecture, Efficiency, and Power Quality',
   sz=26, bold=False, color=RGBColor(0xDD,0xDD,0xDD), align=PP_ALIGN.CENTER)

tb(Inches(1.5), Inches(1.28), Inches(33), Inches(0.38),
   'Jane Yang (yjane@stanford.edu)    Liang Min (liangmin@stanford.edu)',
   sz=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tb(Inches(1.5), Inches(1.62), Inches(33), Inches(0.3),
   'Stanford University    CeraWeek 2026',
   sz=22, bold=False, color=RGBColor(0xCC,0xCC,0xCC), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# COLUMN 1: MOTIVATION + ARCHITECTURE
# ═══════════════════════════════════════

# --- Motivation ---
s1h = Inches(7.8)
add_rect(C1, CT, CW, s1h)
section_title(C1, CT, CW, 'ABSTRACT & MOTIVATION')

tf1 = rtb(C1+Inches(0.2), CT+Inches(0.65), CW-Inches(0.4), Inches(6.8))

ap(tf1, 'Why Revisit the Power Path Now?',
   sz=28, bold=True, color=CARDINAL, sa=Pt(8))

ap(tf1, 'AI-training loads are changing the problem from a purely '
   'steady-state efficiency problem into an architecture + power-dynamics '
   'problem. The GPU is the dominant power consumer (~65%), and '
   'synchronized GPU activity creates periodic power swings visible '
   'at the datacenter level.',
   sz=22, color=BLACK, sa=Pt(10))

ap(tf1, 'Three Industry Scenarios:',
   sz=24, bold=True, color=TEAL, sa=Pt(6))

ap(tf1, '1. Traditional AC: 5\u20136 conversion stages '
   '(Xfmr \u2192 UPS \u2192 PDU \u2192 PSU \u2192 VRM). '
   'Mature but inefficient.',
   sz=21, color=BLACK, sa=Pt(5))

ap(tf1, '2. Scenario 2: 69 kV AC \u2192 800 Vdc perimeter conversion. '
   'Fewer stages and a cleaner 800 Vdc downstream path, but the campus '
   'still depends on an AC perimeter conversion boundary.',
   sz=21, color=BLACK, sa=Pt(5))

ap(tf1, '3. Scenario 3: proposed MVDC backbone. A single MV AC/DC front-end '
   'feeds the DC backbone, so the downstream plant stays DC-native and '
   'buffering, storage, and PV can connect without repeated AC re-entry.',
   sz=21, color=BLACK, sa=Pt(8))

# GPU dynamics figure
img(str(FIG_DIR / 'fig1_gpu_power_trace_and_breakdown.png'),
    C1+Inches(0.15), CT+Inches(5.7), width=CW-Inches(0.3))

# Scale/Dynamics/Spectrum
tf_sd = rtb(C1+Inches(0.2), CT+Inches(7.2), CW-Inches(0.4), Inches(0.5))
ap(tf_sd, 'Scale: >100 MW deployments  |  Dynamics: 0.1\u201320 Hz load swings  |  Spectrum: FFT peaks at 0.2\u20133 Hz',
   sz=16, color=DARK_GRAY, sa=Pt(2))

# --- Architecture ---
s2t = CT + s1h + Inches(0.2)
s2h = CH - s1h - Inches(0.2)
add_rect(C1, s2t, CW, s2h)
section_title(C1, s2t, CW, 'PROPOSED ARCHITECTURE', sz=26)

# Expanded 3-scenario text ABOVE the diagram
tf_arch_text = rtb(C1+Inches(0.2), s2t+Inches(0.65), CW-Inches(0.4), Inches(5.5))

ap(tf_arch_text, 'Three Scenarios:',
   sz=24, bold=True, color=TEAL, sa=Pt(6))

ap(tf_arch_text, '1. Traditional AC-centric',
   sz=22, bold=True, color=RGBColor(0x6B,0x6B,0x6B), sa=Pt(3))
ap(tf_arch_text, 'Utility MV AC \u2192 MV/LV Transformer \u2192 UPS (AC\u2192DC\u2192AC) \u2192 PDU \u2192 '
   'Rack PSU (AC\u2192DC) \u2192 VRM (DC\u2192DC) \u2192 GPU. Five to six conversion stages. '
   'Highest maturity but lowest efficiency due to repeated AC-DC-AC round-trips.',
   sz=20, color=BLACK, sa=Pt(8))

ap(tf_arch_text, '2. AC-fed SST / 800 Vdc Pod (NVIDIA et al.)',
   sz=22, bold=True, color=RGBColor(0x15,0x65,0xC0), sa=Pt(3))
ap(tf_arch_text, 'Utility 69 kV AC \u2192 perimeter AC/DC conversion \u2192 800 Vdc bus \u2192 '
   'rack/node DC-DC \u2192 board DC/DC \u2192 GPU. '
   f'{major_conv_s2} major conversion stages in the current model. '
   'This reduces downstream complexity, but the campus still depends on an '
   'AC perimeter conversion boundary rather than a DC backbone.',
   sz=20, color=BLACK, sa=Pt(8))

ap(tf_arch_text, '3. Proposed: MVDC Hub + Isolated DC Pod',
   sz=22, bold=True, color=RGBColor(0x2E,0x7D,0x32), sa=Pt(3))
ap(tf_arch_text, 'Utility MV AC \u2192 Centralized MV AC/DC Front-End (PFC + Harmonics) \u2192 '
   'MVDC Backbone (DC Subtransmission) \u2192 HF Isolated DC Pod \u2192 800 Vdc Bus \u2192 VRM \u2192 GPU. '
   f'{major_conv_s3} major conversion stages in the current model. Harmonics and PF are handled once at the centralized '
   'front-end. The DC backbone natively supports BESS, solar PV, and rack-level buffering.',
   sz=20, color=BLACK, sa=Pt(8))

ap(tf_arch_text, 'Target: Utility AC \u2192 MV AC/DC \u2192 MVDC Hub \u2192 Isolated DC Pod \u2192 GPU',
   sz=22, bold=True, color=TEAL, sa=Pt(4))

ap(tf_arch_text,
   f'\u03b7_total = \u220f(\u03b7_i):  Scenario 1 = {eff_s1:.2f}%  |  Scenario 2 = {eff_s2:.2f}%  |  Scenario 3 = {eff_s3:.2f}%',
   sz=20, bold=True, color=CARDINAL, sa=Pt(2))

# Architecture diagram below the text
arch_diagram_top = s2t + Inches(6.4)
img(str(FIG_DIR / 'fig3_three_scenario_architecture.png'),
    C1+Inches(0.1), arch_diagram_top, width=CW-Inches(0.2))


# ═══════════════════════════════════════
# COLUMN 2: THREE BENEFITS
# ═══════════════════════════════════════

# --- Benefit 1: Efficiency ---
b1h = Inches(7.0)
add_rect(C2, CT, CW, b1h)
section_title(C2, CT, CW, 'BENEFIT 1: CUMULATIVE EFFICIENCY')

tf_e = rtb(C2+Inches(0.2), CT+Inches(0.65), CW-Inches(0.4), Inches(1.6))
ap(tf_e, 'The MVDC architecture eliminates redundant conversions, achieving '
   f'{eff_s3:.2f}% full-load efficiency versus {eff_s1:.2f}% for Scenario 1 '
   f'and {eff_s2:.2f}% for Scenario 2.',
   sz=22, color=BLACK, sa=Pt(6))
ap(tf_e,
   f'At 100 MW IT, Scenario 3 avoids {delta_loss_gwh_s3_vs_s1:.2f} GWh/year of loss '
   f'and about ${delta_cost_musd_s3_vs_s1:.2f}M/year versus Scenario 1.',
   sz=22, bold=True, color=TEAL, sa=Pt(4))

# Efficiency chart
img(str(FIG_DIR / 'fig2_cumulative_efficiency.png'),
    C2+Inches(0.1), CT+Inches(2.8), width=CW-Inches(0.2))

tb(C2+Inches(0.2), CT+Inches(6.3), CW-Inches(0.4), Inches(0.4),
   'Sources: IEEE PELS 2024; EPRI DC Data Center Report; Rasmussen WP#63',
   sz=15, color=DARK_GRAY)

# --- Benefit 2: Harmonics & PQ (text only, no chart) ---
b2t = CT + b1h + Inches(0.2)
b2h = Inches(5.5)
add_rect(C2, b2t, CW, b2h)
section_title(C2, b2t, CW, 'BENEFIT 2: HARMONICS & POWER QUALITY')

tf_h = rtb(C2+Inches(0.2), b2t+Inches(0.65), CW-Inches(0.4), b2h-Inches(0.8))
ap(tf_h, 'In AC architectures, every server PSU injects harmonics and '
   'requires power factor correction (PFC). With thousands of servers, '
   f'Scenario 1 distributes harmonic ownership across {pq_points_s1} AC-side interfaces.',
   sz=22, color=BLACK, sa=Pt(10))

ap(tf_h, 'Scenario 2 and Scenario 3 both reduce the AC-side count to one boundary, '
   'but only Scenario 3 moves that boundary upstream to the backbone and '
   'keeps the downstream plant DC-native.',
   sz=22, bold=True, color=TEAL, sa=Pt(10))

ap(tf_h, 'For utilities: AC-side power quality is managed only once at the '
   'grid interface \u2014 dramatically simplifying compliance with IEEE 519 '
   'harmonic limits and reducing filter equipment cost.',
   sz=22, color=BLACK, sa=Pt(10))

ap(tf_h, f'Scenario 1: {pq_points_s1} AC harmonic-injection points\n'
   f'Scenario 2: {pq_points_s2} AC boundary point\n'
   f'Scenario 3: {pq_points_s3} centralized AC boundary point',
   sz=22, bold=True, color=CARDINAL, sa=Pt(2))

# --- Benefit 3: Voltage & DC-Native (with power capacity chart) ---
b3t = b2t + b2h + Inches(0.2)
b3h = CH - b1h - b2h - Inches(0.4)
add_rect(C2, b3t, CW, b3h)
section_title(C2, b3t, CW, 'BENEFIT 3: VOLTAGE MGMT & DC-NATIVE PATH', sz=24)

tf_v = rtb(C2+Inches(0.2), b3t+Inches(0.65), CW-Inches(0.4), Inches(2.6))
ap(tf_v, 'The MVDC backbone provides a cleaner DC-native path for '
   'rack buffering, battery storage (BESS), solar PV, and future '
   'modular AI power blocks \u2014 all connecting directly to the DC bus '
   'without AC-DC-AC re-conversion.',
   sz=22, color=BLACK, sa=Pt(8))

ap(tf_v, 'DC voltage regulation is inherently simpler: no reactive power, '
   'no frequency control, no phase balancing.',
   sz=22, color=BLACK, sa=Pt(8))

ap(tf_v, 'At the 34.5/35 kV benchmark in the companion figure, '
   'bipolar DC carries about 23% more power on the same conductor set.',
   sz=22, bold=True, color=TEAL, sa=Pt(4))

# Power capacity chart in Benefit 3
img(str(FIG_DIR / 'fig4_power_capacity_ac_vs_dc.png'),
    C2+Inches(0.1), b3t+Inches(3.6), width=CW-Inches(0.2))


# ═══════════════════════════════════════
# COLUMN 3: COMPARISON TABLE + CONCLUSION
# ═══════════════════════════════════════

# --- Comparison Table ---
tbl_h = Inches(13.5)
add_rect(C3, CT, CW, tbl_h)
section_title(C3, CT, CW, 'ARCHITECTURE COMPARISON', sz=28)

# Subtitle
tb(C3+Inches(0.2), CT+Inches(0.6), CW-Inches(0.4), Inches(0.5),
   'The proposed path exchanges downstream simplicity for a tougher MVDC infrastructure problem.',
   sz=18, color=DARK_GRAY)

# Table
tbl_top = CT + Inches(1.2)
tbl_left = C3 + Inches(0.15)
tbl_w = CW - Inches(0.3)

rows, cols = 8, 4
ts = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, Inches(10.5))
table = ts.table

cws = [Inches(2.0), Inches(2.6), Inches(2.6), Inches(3.0)]
for i, w in enumerate(cws): table.columns[i].width = w

headers = ['Axis', 'Scenario 1\nTraditional AC', 'Scenario 2\n69 kV AC ->\n800 Vdc', 'Scenario 3\nMVDC Backbone']
hcolors = [DARK_GRAY, RGBColor(0x4A,0x4A,0x4A), CARDINAL, TEAL]

for j, (h, hc) in enumerate(zip(headers, hcolors)):
    c = table.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = hc
    for p in c.text_frame.paragraphs:
        p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = WHITE
        p.font.name = 'Arial'; p.alignment = PP_ALIGN.CENTER
    c.vertical_anchor = MSO_ANCHOR.MIDDLE

data = [
    ['Facility\nInterface',
     'AC-native utility\n+ facility protection',
     '69 kV AC perimeter\nconverter to 800 Vdc',
     'Centralized MV AC/DC\nfeeding MVDC bus'],
    ['Pod\nComplexity',
     'Repeated downstream\nAC/DC + DC/DC stages',
     'Simpler 800 Vdc path,\nbut perimeter AC/DC remains',
     'Simpler DC pod;\nAC/DC moves upstream'],
    ['DC-Native\nIntegration',
     'Extra interfaces for\nstorage / solar / DC',
     'Cleaner downstream DC,\nbut no MVDC backbone',
     'Best fit: BESS, solar,\nrack DC all native'],
    ['Harmonic / PF\nOwnership',
     'Distributed across\nAC facility + servers',
     'One AC boundary at\nthe perimeter converter',
     'Centralized at\ncommon MV front end'],
    ['Maturity',
     'Highest',
     'Prototype to\nemerging products',
     'Lowest today;\nprotection is key hurdle'],
    ['Main\nBottleneck',
     'Repeated conversions,\ncopper, facility sprawl',
     'Perimeter conversion,\nfront-end hardware cost',
     'MVDC protection,\nfault handling'],
]

for i, rd in enumerate(data):
    for j, ct in enumerate(rd):
        c = table.cell(i+1, j); c.text = ct
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(17); p.font.bold = (j==0)
            p.font.color.rgb = BLACK; p.font.name = 'Arial'
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        if j == 3:
            c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xE8,0xF5,0xE9)
        elif j == 0:
            c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xF5,0xF5,0xF5)
        else:
            c.fill.solid(); c.fill.fore_color.rgb = WHITE

# Bottom row merged
bc = table.cell(7, 0); be = table.cell(7, 3); bc.merge(be)
bc.text = ('The MVDC-hub proposal is strongest as a system-architecture story \u2014 '
           'not a claim that one converter solves the whole facility.')
for p in bc.text_frame.paragraphs:
    p.font.size = Pt(17); p.font.bold = True; p.font.color.rgb = CARDINAL
    p.font.name = 'Arial'; p.alignment = PP_ALIGN.CENTER
bc.fill.solid(); bc.fill.fore_color.rgb = RGBColor(0xFD,0xED,0xEC)

# --- Conclusion & Future Work ---
ct3 = CT + tbl_h + Inches(0.2)
ch3 = CH - tbl_h - Inches(0.2)
add_rect(C3, ct3, CW, ch3)
section_title(C3, ct3, CW, 'CONCLUSION & FUTURE WORK')

tf_c = rtb(C3+Inches(0.2), ct3+Inches(0.65), CW-Inches(0.4), ch3-Inches(0.8))

ap(tf_c, 'Conclusion',
   sz=26, bold=True, color=TEAL, sa=Pt(6))
ap(tf_c, 'The MVDC subtransmission backbone offers a compelling path for '
   f'next-generation AI factories: {eff_s3:.2f}% modeled full-load efficiency, '
   f'{delta_loss_gwh_s3_vs_s1:.2f} GWh/year lower electrical loss than Scenario 1, '
   'centralized power quality management, simpler voltage regulation, and native '
   'DC integration for storage and renewables.',
   sz=22, color=BLACK, sa=Pt(10))

ap(tf_c, 'Future Work:',
   sz=26, bold=True, color=CARDINAL, sa=Pt(6))
ap(tf_c, '\u2022 MVDC protection schemes, fault isolation, and safety standards development',
   sz=22, color=BLACK, sa=Pt(5))
ap(tf_c, '\u2022 Techno-economic analysis at 100\u2013500 MW scale with hardware-in-the-loop validation',
   sz=22, color=BLACK, sa=Pt(5))
ap(tf_c, '\u2022 Optimal buffering placement (GPU, rack, pod, or MVDC hub) and BOM analysis',
   sz=22, color=BLACK, sa=Pt(10))

ap(tf_c, 'References',
   sz=22, bold=True, color=CARDINAL, sa=Pt(4))
ap(tf_c, '[1] IEA, "Energy and AI," 2025   '
   '[2] NVIDIA, "800 V HVDC Architecture," 2025   '
   '[3] PNNL, "EMT Modeling of Large Data Centers," 2026\n'
   '[4] IEEE Std 519-2022   '
   '[5] Segan et al., Sustainable Energy, Grids and Networks, 2025   '
   '[6] OpenDSS / DSS C-API local validation, 2026',
   sz=16, color=DARK_GRAY, sa=Pt(2))


# ═══════════════════════════════════════
# FOOTER
# ═══════════════════════════════════════
add_box(0, H-FTR_H, W, FTR_H, DARK_CARD)
tb(Inches(0.5), H-FTR_H+Inches(0.08), Inches(16), Inches(0.34),
   'Jane Yang (yjane@stanford.edu)    Liang Min (liangmin@stanford.edu)    Stanford University',
   sz=20, bold=True, color=WHITE)
tb(Inches(22), H-FTR_H+Inches(0.08), Inches(13.5), Inches(0.34),
   'CeraWeek 2026    DC Subtransmission Backbone for AI Factories',
   sz=20, color=RGBColor(0xCC,0xCC,0xCC), align=PP_ALIGN.RIGHT)

# ── Save ──
out = str(OUTPUT_PATH)
prs.save(out)
print(f"Poster v3 saved: {out} ({os.path.getsize(out)/1024:.0f} KB)")
