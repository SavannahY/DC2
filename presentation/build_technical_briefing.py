#!/usr/bin/env python3
from __future__ import annotations

import json
import math
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = ROOT / "presentation"
FIG_DIR = PRESENTATION_DIR / "figures"
OUTPUT_PPTX = PRESENTATION_DIR / "dc_subtransmission_backbone_technical_briefing.pptx"

WHITEPAPER_FIG_DIR = ROOT / "whitepaper" / "figures"
GITHUB_URL = "https://github.com/SavannahY/DC2"


TITLE_COLOR = RGBColor(16, 49, 99)
ACCENT_BLUE = "#1f77b4"
ACCENT_GREEN = "#2ca02c"
ACCENT_ORANGE = "#ff7f0e"
ACCENT_RED = "#d62728"
BODY_COLOR = RGBColor(40, 40, 40)
MUTED_COLOR = RGBColor(90, 90, 90)


def load_json(name: str) -> dict:
    return json.loads((ROOT / name).read_text(encoding="utf-8"))


def ensure_dirs() -> None:
    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    FIG_DIR.mkdir(parents=True, exist_ok=True)


def add_textbox(slide, left, top, width, height, text, font_size=20, bold=False,
                color=BODY_COLOR, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, bullets, font_size=22, color=BODY_COLOR):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        if not p.runs:
            run = p.add_run()
        else:
            run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return box


def add_slide_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(
        slide,
        Inches(0.6),
        Inches(0.3),
        Inches(12.0),
        Inches(0.5),
        title,
        font_size=28,
        bold=True,
        color=TITLE_COLOR,
    )
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.6),
        Inches(0.85),
        Inches(12.1),
        Inches(0.05),
    ).fill.solid()
    shape = slide.shapes[-1]
    shape.fill.fore_color.rgb = RGBColor(41, 98, 255)
    shape.line.fill.background()
    if subtitle:
        add_textbox(
            slide,
            Inches(0.6),
            Inches(0.95),
            Inches(12.0),
            Inches(0.35),
            subtitle,
            font_size=12,
            color=MUTED_COLOR,
        )


def build_public_proof_chart() -> Path:
    benefit = load_json("public_benefit_report.json")
    harmonic = load_json("public_harmonic_frequency_sweep_report.json")

    location_rows = benefit["benefits"]["location_robustness"]["rows"]
    areas = [row["label"] for row in location_rows]
    s2_v = [row["scenario2m_base_max_poi_drop_pct_points"] for row in location_rows]
    s3_v = [row["scenario3m_base_max_poi_drop_pct_points"] for row in location_rows]

    families = harmonic["families"]
    family_order = ["low_order_dominant", "balanced_filtered", "higher_order_filtered"]
    family_labels = ["Low-order", "Balanced", "Higher-order"]
    s2_h = [families[key]["scenario_2m"]["max_poi_thdv_proxy"] * 100.0 for key in family_order]
    s3_h = [families[key]["scenario_3m"]["max_poi_thdv_proxy"] * 100.0 for key in family_order]

    fig, axes = plt.subplots(1, 2, figsize=(12, 4.8))

    x = np.arange(len(family_labels))
    width = 0.35
    axes[0].bar(x - width / 2, s2_h, width, label="Scenario 2(M)", color=ACCENT_BLUE)
    axes[0].bar(x + width / 2, s3_h, width, label="Scenario 3(M)", color=ACCENT_GREEN)
    axes[0].set_yscale("log")
    axes[0].set_xticks(x)
    axes[0].set_xticklabels(family_labels)
    axes[0].set_ylabel("Max POI THDv proxy (%)")
    axes[0].set_title("SMART-DS same-feeder equal-total harmonic benchmark")
    axes[0].grid(True, axis="y", alpha=0.25, which="both")
    axes[0].legend(frameon=False, fontsize=9)

    x = np.arange(len(areas))
    axes[1].bar(x - width / 2, s2_v, width, label="Scenario 2(M)", color=ACCENT_ORANGE)
    axes[1].bar(x + width / 2, s3_v, width, label="Scenario 3(M)", color=ACCENT_GREEN)
    axes[1].set_xticks(x)
    axes[1].set_xticklabels(areas)
    axes[1].set_ylabel("Max POI voltage-drop proxy (pct-pts)")
    axes[1].set_title("RTS-GMLC mirrored-area voltage benchmark")
    axes[1].grid(True, axis="y", alpha=0.25)

    fig.tight_layout()
    output = FIG_DIR / "public_proof_chart.png"
    fig.savefig(output, dpi=220, bbox_inches="tight")
    plt.close(fig)
    return output


def build_rms_chart() -> Path:
    report = load_json("public_rms_dynamic_report.json")
    rows = [
        row
        for row in report["comparisons"]
        if math.isclose(row["campus_share_of_system_load"], 0.25)
        and row["burst_case_name"] == "mit_ai_burst_900s_p95"
        and row["s2_tds_converged"]
        and row["s3_tds_converged"]
    ]
    rows.sort(key=lambda item: item["grid_mode"])
    labels = ["Normal", "Weakened corridor"]
    s2_v = [rows[0]["s2_min_campus_vpu"], rows[1]["s2_min_campus_vpu"]]
    s3_v = [rows[0]["s3_min_campus_vpu"], rows[1]["s3_min_campus_vpu"]]
    s2_f = [rows[0]["s2_max_freq_dev_hz"], rows[1]["s2_max_freq_dev_hz"]]
    s3_f = [rows[0]["s3_max_freq_dev_hz"], rows[1]["s3_max_freq_dev_hz"]]

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.6))
    x = np.arange(len(labels))
    width = 0.35

    axes[0].bar(x - width / 2, s2_v, width, label="Scenario 2(M)", color=ACCENT_BLUE)
    axes[0].bar(x + width / 2, s3_v, width, label="Scenario 3(M)", color=ACCENT_GREEN)
    axes[0].set_ylim(0.65, 1.05)
    axes[0].set_xticks(x)
    axes[0].set_xticklabels(labels)
    axes[0].set_ylabel("Minimum campus voltage (pu)")
    axes[0].set_title("ANDES RMS benchmark at 25% campus share")
    axes[0].grid(True, axis="y", alpha=0.25)
    axes[0].legend(frameon=False, fontsize=9)

    axes[1].bar(x - width / 2, s2_f, width, label="Scenario 2(M)", color=ACCENT_ORANGE)
    axes[1].bar(x + width / 2, s3_f, width, label="Scenario 3(M)", color=ACCENT_GREEN)
    axes[1].set_xticks(x)
    axes[1].set_xticklabels(labels)
    axes[1].set_ylabel("Max |df| (Hz)")
    axes[1].set_title("MIT-derived 900 s p95 burst frequency response")
    axes[1].grid(True, axis="y", alpha=0.25)

    fig.tight_layout()
    output = FIG_DIR / "rms_benchmark_chart.png"
    fig.savefig(output, dpi=220, bbox_inches="tight")
    plt.close(fig)
    return output


def build_fault_chart() -> Path:
    report = load_json("public_fault_envelope_report.json")
    rows = report["baseline_rows"]

    fig, ax = plt.subplots(figsize=(6.4, 4.6))
    times_ms = np.linspace(0.0, 20.0, 400)
    times_s = times_ms / 1000.0

    colors = [ACCENT_RED, ACCENT_BLUE]
    for row, color in zip(rows, colors):
        r_tot = row["total_r_ohm"]
        l_tot = row["total_l_h"]
        i0 = row["prefault_current_a"]
        i_inf = (report["scenario3m_base_case"]["source_voltage_kv"] * 1000.0) / r_tot
        tau = l_tot / r_tot
        current_ka = (i_inf - (i_inf - i0) * np.exp(-times_s / tau)) / 1000.0
        label = row["location_name"].replace("_", " ")
        ax.plot(times_ms, current_ka, linewidth=2.3, label=label, color=color)

    for marker in (2.0, 5.0, 10.0):
        ax.axvline(marker, linestyle="--", linewidth=1.0, color="gray", alpha=0.5)

    ax.set_xlabel("Clearing time (ms)")
    ax.set_ylabel("Fault current (kA)")
    ax.set_title("Scenario 3(M) reduced-order fault envelope\n100 mH reactor, 0.01 ohm fault resistance")
    ax.grid(True, alpha=0.25)
    ax.legend(frameon=False, fontsize=9)

    output = FIG_DIR / "fault_envelope_chart.png"
    fig.tight_layout()
    fig.savefig(output, dpi=220, bbox_inches="tight")
    plt.close(fig)
    return output


def add_picture(slide, path: Path, left, top, width=None, height=None):
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def build_presentation() -> Path:
    ensure_dirs()

    public_chart = build_public_proof_chart()
    rms_chart = build_rms_chart()
    fault_chart = build_fault_chart()

    source_report = load_json("source_backed_model_report.json")
    campus_report = load_json("multinode_campus_report.json")
    benefit_report = load_json("public_benefit_report.json")
    common_td = load_json("public_common_network_td_report.json")
    rms_report = load_json("public_rms_dynamic_report.json")
    fault_report = load_json("public_fault_envelope_report.json")

    single_results = {row["display_name"]: row for row in source_report["results"]}
    multi_results = {row["name"]: row for row in campus_report["architectures"]}

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_textbox(slide, Inches(0.7), Inches(0.8), Inches(11.8), Inches(0.8),
                "Direct Current Subtransmission Backbone for AI Factories",
                font_size=30, bold=True, color=TITLE_COLOR)
    add_textbox(slide, Inches(0.7), Inches(1.7), Inches(11.6), Inches(0.6),
                "Technical briefing based on the public-data modeling stack in SavannahY/DC2",
                font_size=18, color=MUTED_COLOR)
    add_bullets(
        slide,
        Inches(0.9),
        Inches(2.5),
        Inches(7.2),
        Inches(2.5),
        [
            "Scenario 2 is the explicit AC-fed SST + 800 VDC baseline.",
            "Scenario 3 adds the MVDC subtransmission backbone while keeping the same downstream 800 VDC structure.",
            "Scenario 3(M) is the campus-scale proof layer with a shared 69 kV DC backbone and four compute blocks.",
        ],
        font_size=22,
    )
    add_textbox(slide, Inches(0.9), Inches(6.5), Inches(6.0), Inches(0.3),
                GITHUB_URL, font_size=12, color=MUTED_COLOR)
    add_picture(slide, WHITEPAPER_FIG_DIR / "architecture_comparison.png", Inches(7.7), Inches(1.3), width=Inches(4.8))

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Architecture Under Test", "Scenario 2 and Scenario 3 share the SST + 800 VDC downstream structure; the backbone is the differentiator.")
    add_picture(slide, WHITEPAPER_FIG_DIR / "architecture_comparison.png", Inches(0.8), Inches(1.3), width=Inches(11.8))
    add_textbox(slide, Inches(1.0), Inches(6.7), Inches(11.2), Inches(0.35),
                "Scenario 3(M) extends Scenario 3 into a shared multi-block campus backbone rather than a single equivalent path.",
                font_size=15, color=MUTED_COLOR, align=PP_ALIGN.CENTER)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Evidence Stack", "What is modeled with public data today, and what is still outside the evidence boundary.")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.3),
        Inches(5.7),
        Inches(4.8),
        [
            "Source-anchored path model for Scenarios 1/2/3",
            "Shared-campus multi-node model for Scenarios 1(M)/2(M)/3(M)",
            "Public ESIF annual bins and MIT Supercloud burst library",
            "OpenDSS AC-boundary cross-checks and harmonics-mode scans",
            "RTS-GMLC mirrored-area and local N-1 sensitivity screens",
            "SMART-DS same-feeder common-network and equal-total harmonic benchmark",
            "ANDES RMS dynamic benchmark under normal and weakened corridor stress",
            "Reduced-order MVDC fault-duty and interruption-timescale screen",
        ],
        font_size=19,
    )
    add_bullets(
        slide,
        Inches(7.0),
        Inches(1.5),
        Inches(5.3),
        Inches(4.4),
        [
            "Supported now:",
            "Campus-scale efficiency direction",
            "Centralized AC-boundary PQ and voltage benefits",
            "Weak-grid RMS voltage robustness",
            "Fault interruption timescale as a first-order issue",
            "",
            "Not proven yet:",
            "IEEE 519 compliance margins",
            "Converter EMT stability and controls",
            "Deployable selective protection design",
            "Site-specific interconnection margins",
        ],
        font_size=19,
    )

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Steady-State and Campus-Scale Results")
    add_picture(slide, WHITEPAPER_FIG_DIR / "cumulative_efficiency.png", Inches(0.7), Inches(1.3), width=Inches(6.4))
    s2 = single_results["AC-fed SST + 800 VDC baseline"]
    s3 = single_results["Proposed MVDC backbone"]
    s2m = multi_results["ac_fed_sst_800vdc"]
    s3m = multi_results["proposed_mvdc_backbone"]
    add_bullets(
        slide,
        Inches(7.4),
        Inches(1.4),
        Inches(5.2),
        Inches(4.8),
        [
            f"Single-path, 100 MW IT: S2 {s2['full_load_total_efficiency']*100:.2f}% vs S3 {s3['full_load_total_efficiency']*100:.2f}%",
            f"Annual loss: {s2['annual_loss_mwh']/1000:.2f} vs {s3['annual_loss_mwh']/1000:.2f} GWh/yr",
            f"Multi-node, 100 MW IT: S2(M) {s2m['full_load']['total_efficiency']*100:.2f}% vs S3(M) {s3m['full_load']['total_efficiency']*100:.2f}%",
            f"Multi-node annual loss: {s2m['annual_summary']['annual_loss_mwh']/1000:.2f} vs {s3m['annual_summary']['annual_loss_mwh']/1000:.2f} GWh/yr",
            "The single-path gain is modest.",
            "The campus-scale multi-node result is the stronger efficiency case.",
        ],
        font_size=20,
    )

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Public Network Proof of Benefits 2 and 3")
    add_picture(slide, public_chart, Inches(0.7), Inches(1.35), width=Inches(7.2))
    family_ratios = [
        load_json("public_harmonic_frequency_sweep_report.json")["families"][key]["max_poi_thdv_ratio_s2_to_s3"]
        for key in ["low_order_dominant", "balanced_filtered", "higher_order_filtered"]
    ]
    td_s2 = common_td["n_minus_one"]["Scenario 2(M)"]["surviving_count"]
    td_s3 = common_td["n_minus_one"]["Scenario 3(M)"]["surviving_count"]
    add_bullets(
        slide,
        Inches(8.0),
        Inches(1.4),
        Inches(4.7),
        Inches(4.8),
        [
            f"RTS mirrored areas: harmonic proxy ratio >= {min(row['scenario2m_harmonic_thdv_proxy_pu']/row['scenario3m_harmonic_thdv_proxy_pu'] for row in benefit_report['benefits']['location_robustness']['rows']):.2f}x",
            f"RTS mirrored areas: voltage-drop ratio >= {min(row['scenario2m_base_max_poi_drop_pct_points']/row['scenario3m_base_max_poi_drop_pct_points'] for row in benefit_report['benefits']['location_robustness']['rows']):.2f}x",
            f"SMART-DS same-feeder equal-total harmonic benchmark: {min(family_ratios):.2f}x to {max(family_ratios):.2f}x lower THDv proxy",
            f"SMART-DS surviving local N-1 cases: {td_s2} for S2(M) vs {td_s3} for S3(M)",
            "This is the strongest current support for centralized AC-boundary PQ ownership and voltage robustness.",
        ],
        font_size=18,
    )

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "RMS Dynamic Benchmark and Fault Envelope")
    add_picture(slide, rms_chart, Inches(0.7), Inches(1.35), width=Inches(6.0))
    add_picture(slide, fault_chart, Inches(6.95), Inches(1.35), width=Inches(5.8))
    strongest = rms_report["headline"]["strongest_voltage_separation"]
    fault_headline = fault_report["headline"]
    add_bullets(
        slide,
        Inches(0.9),
        Inches(6.0),
        Inches(12.0),
        Inches(0.9),
        [
            f"Strongest converged RMS separation: weakened corridor, 900 s p95 burst, 25% campus share, Vmin {strongest['s2_min_campus_vpu']:.4f} pu vs {strongest['s3_min_campus_vpu']:.4f} pu.",
            f"Fault screen: source-backbone midpoint reaches {fault_headline['worst_hybrid_reference_case']['breaker_rows']['hybrid_reference']['current_a']/1000:.2f} kA at 5 ms in the worst public sweep; fastest time to 5 kA is {fault_headline['fastest_time_to_5ka']['time_to_thresholds_s']['5.0kA']*1000:.2f} ms.",
        ],
        font_size=16,
    )

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "What the Current Evidence Supports")
    add_bullets(
        slide,
        Inches(0.9),
        Inches(1.5),
        Inches(5.7),
        Inches(4.5),
        [
            "Supported at the current evidence level:",
            "Scenario 3(M) is the stronger campus-scale architecture case.",
            "Benefit 2: centralized AC-boundary harmonic sensitivity is lower.",
            "Benefit 3: upstream AC-boundary voltage and weak-grid RMS robustness are better.",
            "Fault interruption timescale is a first-order design issue for the MVDC backbone.",
        ],
        font_size=20,
    )
    add_bullets(
        slide,
        Inches(7.0),
        Inches(1.5),
        Inches(5.4),
        Inches(4.5),
        [
            "Still not proven:",
            "Exact IEEE 519 compliance margins",
            "Converter-aware EMT stability",
            "Deployable breaker and protection settings",
            "Site-specific utility interconnection limits",
            "Uniform internal dynamic stress reduction",
        ],
        font_size=20,
    )

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Next Technical Steps")
    add_bullets(
        slide,
        Inches(0.9),
        Inches(1.5),
        Inches(11.8),
        Inches(4.7),
        [
            "Replace proxy curves for the centralized AC/DC front end and isolated DC pod with measured or vendor-qualified data.",
            "Promote the RMS benchmark into converter-aware EMT studies for the centralized front end, isolated pod, and buffer controls.",
            "Turn the fault envelope into a selective protection and grounding study with explicit interruption technology assumptions.",
            "Extend the public common-network workflow into a site-specific PCC study when utility Thevenin and short-circuit data are available.",
            "Use the repo as the reproducible base package: code, figures, white paper, and slide deck are all versioned together.",
        ],
        font_size=21,
    )
    add_textbox(slide, Inches(0.9), Inches(6.55), Inches(11.8), Inches(0.25),
                GITHUB_URL, font_size=12, color=MUTED_COLOR, align=PP_ALIGN.CENTER)

    prs.save(OUTPUT_PPTX)
    return OUTPUT_PPTX


if __name__ == "__main__":
    path = build_presentation()
    print(f"Wrote {path}")
